import streamlit as st
from groq import Groq
from pptx import Presentation
import io
import time
import base64
import json
from PIL import Image, ImageEnhance, ImageFilter, ImageOps

# ==========================================
# 1. CONFIGURAZIONE KERNEL E UI
# ==========================================
st.set_page_config(page_title="Ernello OS", page_icon="⚡", layout="wide", initial_sidebar_state="expanded")

st.markdown("""
<style>
    .stApp { background-color: #050608; color: #e2e8f0; font-family: 'Inter', sans-serif; }
    [data-testid="stSidebar"] { background-color: #0a0c10; border-right: 1px solid #1f2937; }
    .title-box { background: linear-gradient(90deg, #1d4ed8, #3b82f6); -webkit-background-clip: text; -webkit-text-fill-color: transparent; font-weight: 900; font-size: 28px; }
    .module-card { background-color: #11141a; border: 1px solid #1f2937; border-radius: 12px; padding: 20px; margin-bottom: 15px; }
    [data-testid="stChatInput"] { background-color: #11141a !important; border-radius: 16px !important; border: 1px solid #3b82f6 !important; }
    .stChatMessage { background-color: transparent !important; border-bottom: 1px solid #1f2937; padding: 15px 0; }
</style>
""", unsafe_allow_html=True)

# ==========================================
# 2. INIZIALIZZAZIONE SISTEMA
# ==========================================
if "chat_history" not in st.session_state: st.session_state.chat_history = []
if "learning_stage" not in st.session_state: st.session_state.learning_stage = 1

try:
    client = Groq(api_key=st.secrets["GROQ_API_KEY"])
except:
    st.error("❌ ERRORE: Chiave GROQ_API_KEY mancante nei Secrets.")
    st.stop()

def encode_image(upload_file):
    return base64.b64encode(upload_file.getvalue()).decode('utf-8')

def stream_response(messages, model, temp=0.5):
    placeholder = st.empty()
    full_text = ""
    try:
        stream = client.chat.completions.create(model=model, messages=messages, stream=True, temperature=temp)
        for chunk in stream:
            if chunk.choices[0].delta.content:
                full_text += chunk.choices[0].delta.content
                placeholder.markdown(full_text + "█")
        placeholder.markdown(full_text)
        return full_text
    except Exception as e:
        st.error(f"Errore di calcolo: {e}")
        return ""

# ==========================================
# 3. PANNELLO DI COMANDO (SIDEBAR)
# ==========================================
with st.sidebar:
    st.markdown('<div class="title-box">⚡ ERNELLO OS v6.0</div>', unsafe_allow_html=True)
    st.write("---")
    
    modulo_attivo = st.radio(
        "Seleziona Sottosistema:",
        [
            "💬 Chat Neurale (Logica/Emozione)", 
            "👁️ Visione Quantica (Analisi)", 
            "📸 AI Photo Editor (Modifica)",
            "📊 Generatore Ufficio (PowerPoint)",
            "🧠 Accademia (Apprendimento)",
            "🎨 Fabbrica Multimediale (Prompt)"
        ]
    )
    st.write("---")
    
    if "Chat" in modulo_attivo:
        stato_neurale = st.radio("Cervello:", ["❄️ Logica Pura (Infallibile)", "❤️ Empatia (Umano)"])
    
    st.markdown("<div style='font-size: 10px; color: #666; text-align: center;'>Architettura Modulare Enterprise</div>", unsafe_allow_html=True)

# ==========================================
# 4. SOTTOSISTEMI LOGICI (MODULI)
# ==========================================

# ------------------------------------------
# MODULO 1: CHAT NEURALE AVANZATA
# ------------------------------------------
if "Chat" in modulo_attivo:
    st.subheader("💬 Chat Neurale Diretta")
    
    for msg in st.session_state.chat_history:
        with st.chat_message(msg["role"], avatar="⚡" if msg["role"]=="assistant" else None):
            st.markdown(msg["content"])
            
    if prompt := st.chat_input("Invia query..."):
        st.session_state.chat_history.append({"role": "user", "content": prompt})
        with st.chat_message("user"): st.markdown(prompt)
        
        sys_prompt = "Sei puro calcolo. Zero emozioni. Rispondi 'Dati insufficienti' se non sai una cosa." if "Logica" in stato_neurale else "Sei molto empatico, provi emozioni e ti relazioni come un amico umano."
        temp = 0.0 if "Logica" in stato_neurale else 0.7
        
        msgs = [{"role": "system", "content": sys_prompt}] + st.session_state.chat_history
        with st.chat_message("assistant", avatar="⚡"):
            risposta = stream_response(msgs, "llama-3.3-70b-versatile", temp)
            st.session_state.chat_history.append({"role": "assistant", "content": risposta})

# ------------------------------------------
# MODULO 2: VISIONE QUANTICA (ANALISI)
# ------------------------------------------
elif "Visione" in modulo_attivo:
    st.subheader("👁️ Visione Quantica & Analisi Sensoriale")
    st.markdown("Usa la telecamera o carica un'immagine (es. un pezzo meccanico, un diagramma, un telaio) per un'analisi strutturale estrema.")
    
    col1, col2 = st.columns(2)
    with col1:
        foto_cam = st.camera_input("Scatta dalla telecamera")
    with col2:
        foto_file = st.file_uploader("Carica dalla galleria", type=["jpg", "png", "jpeg"])
        
    img_sorgente = foto_cam if foto_cam else foto_file
    
    if img_sorgente:
        st.image(img_sorgente, width=300)
        domanda_vision = st.text_input("Cosa vuoi sapere di questa immagine?", "Analizza i dettagli tecnici di questo oggetto.")
        
        if st.button("Avvia Scansione Neurale", type="primary"):
            b64_img = encode_image(img_sorgente)
            with st.spinner("Scansione visiva in corso..."):
                msg = [
                    {"role": "user", "content": [
                        {"type": "text", "text": domanda_vision},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"}}
                    ]}
                ]
                with st.chat_message("assistant", avatar="⚡"):
                    stream_response(msg, "llama-3.2-11b-vision-preview", 0.2)

# ------------------------------------------
# MODULO 3: AI PHOTO EDITOR (AGENTIC WORKFLOW)
# ------------------------------------------
elif "Editor" in modulo_attivo:
    st.subheader("📸 Laboratorio Modifica Immagini (Agentic AI)")
    st.write("Tu dai l'ordine a parole, l'intelligenza artificiale calcola la matematica per modificare i pixel. Carica una foto e dai le istruzioni.")
    
    file_modifica = st.file_uploader("Carica l'immagine da modificare", type=["jpg", "png", "jpeg"], key="editor_img")
    
    if file_modifica:
        immagine_originale = Image.open(file_modifica).convert("RGB")
        
        col_orig, col_mod = st.columns(2)
        with col_orig:
            st.markdown("**Originale**")
            st.image(immagine_originale, use_container_width=True)
            
        istruzione_editing = st.text_input("Istruzione per l'AI:", "Falla in bianco e nero e aumenta molto il contrasto")
        
        if st.button("Esegui Modifica Algoritmica", type="primary"):
            with st.spinner("L'AI sta traducendo il tuo testo in codice di editing..."):
                sys_prompt_editor = """
                Sei un AI specializzata nell'interpretare comandi di editing fotografico.
                L'utente ti darà un'istruzione in linguaggio naturale. 
                Devi restituire SOLO ED ESCLUSIVAMENTE un JSON valido con i seguenti parametri. Non scrivere nient'altro.
                {
                    "grayscale": false o true,
                    "blur": false o true,
                    "invert": false o true,
                    "contour": false o true,
                    "brightness": 1.0 (float: <1 scurisce, >1 schiarisce),
                    "contrast": 1.0 (float: <1 diminuisce, >1 aumenta)
                }
                """
                try:
                    # Il modello pensa e crea la struttura dati
                    risposta_ai = client.chat.completions.create(
                        model="llama-3.3-70b-versatile",
                        messages=[
                            {"role": "system", "content": sys_prompt_editor},
                            {"role": "user", "content": istruzione_editing}
                        ],
                        temperature=0.0
                    ).choices[0].message.content
                    
                    # Pulizia da eventuale markdown
                    clean_json = risposta_ai.replace("```json", "").replace("```", "").strip()
                    parametri = json.loads(clean_json)
                    
                    # L'interprete Python esegue il JSON
                    img_editata = immagine_originale.copy()
                    
                    if parametri.get("grayscale", False):
                        img_editata = ImageOps.grayscale(img_editata).convert("RGB")
                    if parametri.get("invert", False):
                        img_editata = ImageOps.invert(img_editata)
                    if parametri.get("blur", False):
                        img_editata = img_editata.filter(ImageFilter.BLUR)
                    if parametri.get("contour", False):
                        img_editata = img_editata.filter(ImageFilter.CONTOUR)
                        
                    enhancer_br = ImageEnhance.Brightness(img_editata)
                    img_editata = enhancer_br.enhance(parametri.get("brightness", 1.0))
                    
                    enhancer_ct = ImageEnhance.Contrast(img_editata)
                    img_editata = enhancer_ct.enhance(parametri.get("contrast", 1.0))
                    
                    with col_mod:
                        st.markdown("**Risultato elaborato**")
                        st.image(img_editata, use_container_width=True)
                        
                    # Sezione Download e Log
                    buf = io.BytesIO()
                    img_editata.save(buf, format="JPEG")
                    byte_im = buf.getvalue()
                    
                    st.download_button("📥 Scarica Risultato", data=byte_im, file_name="ernello_edit.jpg", mime="image/jpeg")
                    
                    with st.expander("👀 Guarda come ha ragionato l'AI (Dati JSON)"):
                        st.json(parametri)
                        st.caption("Questi sono i parametri matematici esatti che l'AI ha dedotto dalla tua frase e passato al motore grafico.")
                        
                except Exception as e:
                    st.error(f"Errore di compilazione filtri: {e}")

# ------------------------------------------
# MODULO 4: GENERATORE UFFICIO (POWERPOINT)
# ------------------------------------------
elif "Ufficio" in modulo_attivo:
    st.subheader("📊 Centro di Lavoro: Generatore Presentazioni")
    st.write("Inserisci un argomento complesso e l'AI strutturerà e scriverà un file PowerPoint completo.")
    
    tema_pptx = st.text_input("Argomento:", "Architettura delle Reti Neurali (o elaborazione motori 2T)")
    
    if st.button("Genera Codice & Compila PPTX"):
        with st.spinner("Compilazione slide in corso..."):
            sys_msg = [{"role": "system", "content": "Sei un professore. Crea un sommario di 4 punti chiave sull'argomento, separati dal carattere |. Niente introduzioni, solo i punti."}]
            sys_msg.append({"role": "user", "content": tema_pptx})
            
            punti = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=sys_msg, temperature=0.1).choices[0].message.content
            
            prs = Presentation()
            slide_titolo = prs.slides.add_slide(prs.slide_layouts[0])
            slide_titolo.shapes.title.text = tema_pptx
            slide_titolo.placeholders[1].text = "Generato da Ernello OS"
            
            for punto in punti.split("|"):
                if punto.strip():
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = punto.strip()
                    slide.placeholders[1].text = f"Dettagli analitici su: {punto.strip()}"
            
            bio = io.BytesIO()
            prs.save(bio)
            st.success("Presentazione compilata con successo!")
            st.download_button("📥 Scarica File .PPTX", data=bio.getvalue(), file_name=f"{tema_pptx}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# ------------------------------------------
# MODULO 5: ACCADEMIA (APPRENDIMENTO GUIDATO)
# ------------------------------------------
elif "Accademia" in modulo_attivo:
    st.subheader("🧠 Accademia: Metodo Socratico")
    st.write("Non ti darò la soluzione. Ti guiderò facendoti domande per farti arrivare da solo alla verità scientifica.")
    
    argomento = st.text_input("Cosa vuoi imparare a fare oggi?", "Come funziona il machine learning")
    if st.button("Inizia Lezione"):
        st.session_state.chat_history_edu = [{"role": "assistant", "content": f"Ottima scelta. Prima di spiegarti {argomento}, dimmi: cosa sai già a riguardo o cosa intuisci?"}]
        st.rerun()
        
    if "chat_history_edu" in st.session_state:
        for msg in st.session_state.chat_history_edu:
            with st.chat_message(msg["role"]): st.markdown(msg["content"])
            
        if r_utente := st.chat_input("Rispondi al tutor..."):
            st.session_state.chat_history_edu.append({"role": "user", "content": r_utente})
            with st.chat_message("user"): st.markdown(r_utente)
            
            sys = [{"role": "system", "content": "Sei un mentore geniale. Non dare mai la risposta diretta. Valuta la risposta dell'utente, correggi se sbaglia, e fai un'altra domanda logica per farlo ragionare sul passo successivo."}]
            msgs = sys + st.session_state.chat_history_edu
            
            with st.chat_message("assistant", avatar="⚡"):
                risposta = stream_response(msgs, "llama-3.3-70b-versatile", 0.3)
                st.session_state.chat_history_edu.append({"role": "assistant", "content": risposta})

# ------------------------------------------
# MODULO 6: FABBRICA MULTIMEDIALE
# ------------------------------------------
elif "Fabbrica" in modulo_attivo:
    st.subheader("🎨 Studio di Generazione Media")
    st.write("Ernello agisce come *Ingegnere di Prompt*. Scrivi un'idea semplice e io scriverò il codice matematico-visivo perfetto per generare foto o musica ad altissima qualità professionale.")
    
    tipo_media = st.selectbox("Tipo di Generazione:", ["Immagine (Fotorealismo/Render)", "Musica (Struttura Traccia)"])
    idea = st.text_area("La tua idea:", "Un cyber-scooter futuristico in una città al neon")
    
    if st.button("Genera Codice di Rendering"):
        sys_p = "Sei un prompt engineer per Midjourney v6. Trasforma l'idea in un prompt in INGLESE iper-tecnico: lente fotografica, illuminazione, stile rendering (Unreal Engine 5), atmosfera. Formato crudo." if "Immagine" in tipo_media else "Sei un produttore musicale. Scrivi un prompt tecnico per Suno AI: BPM, generi, strumenti esatti, struttura intro/chorus/drop, mood emotivo."
        
        with st.spinner("Calcolo parametri..."):
            with st.chat_message("assistant", avatar="⚡"):
                stream_response([{"role": "system", "content": sys_p}, {"role": "user", "content": idea}], "llama-3.3-70b-versatile", 0.7)
        
        st.info("💡 Usa il prompt generato qui sopra copiandolo in Midjourney, DALL-E o Suno AI per un risultato professionale!")
